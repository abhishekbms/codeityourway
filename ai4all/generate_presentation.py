"""
AI 4 All — Day 1 Session: Agent Flows in Microsoft Copilot Studio
Generates AI4All_Day1_AgentFlows.pptx with 14 fully formatted slides.
Source: https://learn.microsoft.com/en-us/microsoft-copilot-studio/flows-overview
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour constants ────────────────────────────────────────────────────────
MS_BLUE = RGBColor(0x00, 0x78, 0xD4)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x20, 0x20, 0x20)
GREY    = RGBColor(0x70, 0x70, 0x70)
LIGHT_BLUE = RGBColor(0xF3, 0xF9, 0xFF)

FOOTER_TEXT = "AI 4 All — Day 1 | Agent Flows in Microsoft Copilot Studio | March 26, 2026"

# ── Slide dimensions (16:9 widescreen) ─────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ─────────────────────────────────────────────────────────────────

def set_cell_bg(cell, rgb_color: RGBColor):
    """Fill a table cell background with a solid colour."""
    from pptx.oxml.ns import qn
    from lxml import etree
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", "{:02X}{:02X}{:02X}".format(rgb_color[0], rgb_color[1], rgb_color[2]))


def add_textbox(slide, left, top, width, height, text, font_size=16,
                bold=False, color=DARK, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_title(slide, title_text, top=Inches(0.35), color=MS_BLUE, font_size=28):
    left  = Inches(0.5)
    width = Inches(12.33)
    height = Inches(0.7)
    tb = add_textbox(slide, left, top, width, height, title_text,
                     font_size=font_size, bold=True, color=color)
    return tb


def add_footer_and_slide_number(slide, slide_number):
    """Add footer text (bottom-left) and slide number (bottom-right)."""
    footer_top = Inches(7.1)
    footer_left = Inches(0.4)
    footer_width = Inches(11.5)
    footer_height = Inches(0.35)
    add_textbox(slide, footer_left, footer_top, footer_width, footer_height,
                FOOTER_TEXT, font_size=9, color=GREY)

    num_left = Inches(12.5)
    num_width = Inches(0.7)
    add_textbox(slide, num_left, footer_top, num_width, footer_height,
                str(slide_number), font_size=9, color=GREY, align=PP_ALIGN.RIGHT)


def add_body_textbox(slide, left, top, width, height):
    """Return a textbox whose text_frame we can fill manually."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_paragraph(tf, text, level=0, font_size=None, bold=False, color=DARK):
    """Append a paragraph to a text_frame with proper indentation."""
    p = tf.add_paragraph()
    p.level = level
    p.space_before = Pt(2)
    if level == 0:
        p.space_before = Pt(4)

    run = p.add_run()
    run.text = text

    if font_size is None:
        font_size = 13 if level > 0 else 16
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


# ── Slide builders ───────────────────────────────────────────────────────────

def build_slide1(prs):
    """Title slide — full-bleed blue background, white text."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Blue background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = MS_BLUE

    # Main title
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(11.9), Inches(1.2),
                "AI 4 ALL — Day 1 Session",
                font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, Inches(0.7), Inches(2.8), Inches(11.9), Inches(0.7),
                "Agent Flows in Microsoft Copilot Studio",
                font_size=26, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    # Body lines
    body_lines = [
        "From Automation to Autonomous Intelligence",
        "Presenter: abhishekbmsI",
        "Date: March 26, 2026",
        "Source: learn.microsoft.com/microsoft-copilot-studio/flows-overview",
    ]
    top = Inches(3.7)
    for line in body_lines:
        add_textbox(slide, Inches(0.7), top, Inches(11.9), Inches(0.4),
                    line, font_size=16, color=WHITE, align=PP_ALIGN.CENTER)
        top += Inches(0.4)

    return slide


def build_slide2(prs):
    """Agenda slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_title(slide, "📌 Today's Agenda")

    tf = add_body_textbox(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.6))
    tf.paragraphs[0].runs  # just reference to not leave an empty first para

    bullets = [
        ("1.  What is Copilot Studio? [5 min]", 0),
        ("2.  What are Agent Flows? [10 min]", 0),
        ("3.  Agent Flows vs Cloud Flows [10 min]", 0),
        ("4.  Architecture: Triggers, Actions, Connectors [10 min]", 0),
        ("5.  Creating an Agent Flow (Step-by-Step) [15 min]", 0),
        ("6.  Using Flows in Agents (as Tools) [10 min]", 0),
        ("7.  Autonomous Agents & AI Capabilities [10 min]", 0),
        ("8.  🧪 LIVE DEMO: Build a Help Desk Agent Flow [30 min]", 0),
        ("9.  Q&A + Resources [10 min]", 0),
    ]

    first = True
    for text, level in bullets:
        if first:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(16)
            run.font.color.rgb = DARK
            first = False
        else:
            add_paragraph(tf, text, level=level, font_size=16)

    add_footer_and_slide_number(slide, 2)
    return slide


def build_slide3(prs):
    """What is Microsoft Copilot Studio?"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_title(slide, "🏢 Microsoft Copilot Studio")

    add_textbox(slide, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.45),
                '"The Low-Code Platform for Building AI Agents"',
                font_size=16, bold=True, color=MS_BLUE)

    tf = add_body_textbox(slide, Inches(0.5), Inches(1.75), Inches(12.3), Inches(5.0))
    bullets = [
        ("Build custom AI agents without deep coding", 0),
        ("Integrate with 1000+ connectors via Power Automate", 0),
        ("Supports both conversational & autonomous agents", 0),
        ("Built on Microsoft Azure AI foundations", 0),
        ("Combines: Conversational AI (Topics) + Agent Flows + Power Platform", 0),
    ]
    first = True
    for text, level in bullets:
        if first:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(16)
            run.font.color.rgb = DARK
            first = False
        else:
            add_paragraph(tf, text, level=level, font_size=16)

    add_footer_and_slide_number(slide, 3)
    return slide


def build_slide4(prs):
    """What Are Agent Flows?"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_title(slide, "🌊 What Are Agent Flows?")

    tf = add_body_textbox(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.9))
    bullets = [
        ("Automated workflows built NATIVELY in Copilot Studio", 0),
        ("Powered by the Power Automate engine", 0),
        ("Pattern: TRIGGER → ACTIONS → CONDITIONS → OUTPUT", 0),
        ("✅ Deterministic — same input = same output", 0),
        ("✅ Low-code/No-code — visual drag-and-drop designer", 0),
        ("✅ Natural Language creation — describe and it builds", 0),
        ("✅ AI-first — built-in AI actions available", 0),
        ("✅ Billed via Copilot Studio — no extra PA license needed", 0),
        ('💡 They are the "hands" of your AI agent — they DO things on behalf of the agent', 0),
    ]
    first = True
    for text, level in bullets:
        if first:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(16)
            run.font.color.rgb = DARK
            first = False
        else:
            add_paragraph(tf, text, level=level, font_size=16)

    add_footer_and_slide_number(slide, 4)
    return slide


def build_slide5(prs):
    """Agent Flows vs Cloud Flows — comparison table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_title(slide, "⚡ Agent Flows vs ☁️ Cloud Flows")

    rows = [
        ["Feature", "Agent Flows (Copilot Studio)", "Cloud Flows (Power Automate)"],
        ["Home", "Lives in Copilot Studio", "Lives in Power Automate portal"],
        ["Purpose", "Agent-centric automation", "General-purpose automation"],
        ["Billing", "Copilot Studio consumption", "Requires Power Automate license"],
        ["Collaboration", "No sharing/co-owners", "Share, co-own, run-only permissions"],
        ["Connectors", "All Power Platform connectors", "All Power Platform connectors"],
        ["Agent Use", "Used directly as agent tools", "Can be called from agents via action step"],
    ]

    col_widths = [Inches(2.2), Inches(4.8), Inches(4.8)]
    row_height = Inches(0.55)
    left = Inches(0.7)
    top  = Inches(1.25)

    tbl = slide.shapes.add_table(
        len(rows), len(rows[0]), left, top,
        sum(col_widths), row_height * len(rows)
    ).table

    # Set column widths
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    for ri, row_data in enumerate(rows):
        for ci, cell_text in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            cell.text = cell_text
            tf = cell.text_frame
            tf.paragraphs[0].runs[0].font.size = Pt(13)
            tf.paragraphs[0].runs[0].font.bold = (ri == 0)

            if ri == 0:
                # Header row
                tf.paragraphs[0].runs[0].font.color.rgb = WHITE
                set_cell_bg(cell, MS_BLUE)
            elif ri % 2 == 0:
                set_cell_bg(cell, LIGHT_BLUE)
                tf.paragraphs[0].runs[0].font.color.rgb = DARK
            else:
                tf.paragraphs[0].runs[0].font.color.rgb = DARK

    add_footer_and_slide_number(slide, 5)
    return slide


def build_slide6(prs):
    """The 3-Layer Architecture."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_title(slide, "🏗️ Agent Flow Architecture")

    tf = add_body_textbox(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.9))

    content = [
        ('LAYER 1: TRIGGERS — "What starts the flow?"', 0, True),
        ("Manual / On-Demand", 1, False),
        ("Automated (Event-based)", 1, False),
        ("Scheduled (Timer)", 1, False),
        ('LAYER 2: ACTIONS — "What does the flow do?"', 0, True),
        ("Send Message, Fetch Data, AI Prompt, Update Record", 1, False),
        ('LAYER 3: CONNECTORS — "How does it connect to systems?"', 0, True),
        ("M365, SharePoint, Dynamics 365, Salesforce, Custom APIs", 1, False),
    ]

    first = True
    for text, level, bold in content:
        fs = 16 if level == 0 else 13
        if first:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(fs)
            run.font.bold = bold
            run.font.color.rgb = MS_BLUE if bold else DARK
            first = False
        else:
            p = add_paragraph(tf, text, level=level, font_size=fs, bold=bold,
                               color=MS_BLUE if (bold and level == 0) else DARK)

    add_footer_and_slide_number(slide, 6)
    return slide


def build_slide7(prs):
    """Triggers Deep Dive."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_title(slide, "🚦 Trigger Types in Agent Flows")

    tf = add_body_textbox(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.9))

    content = [
        ("1️⃣ MANUAL TRIGGER", 0, True),
        ("User or agent explicitly calls the flow", 1, False),
        ("Accepts input parameters from the caller", 1, False),
        ('Use for: "on-demand" agent actions', 1, False),
        ("2️⃣ AUTOMATED TRIGGER", 0, True),
        ("Fires in response to an external event (e.g., new email, new DB record)", 1, False),
        ("Use for: Background automation, monitoring", 1, False),
        ("3️⃣ SCHEDULED TRIGGER", 0, True),
        ("Runs on a defined time schedule (e.g., every day at 9am)", 1, False),
        ("Use for: Reports, digests, daily tasks", 1, False),
        ("⚡ Flows called from AGENTS use Manual triggers with parameters passed from conversation context", 0, False),
    ]

    first = True
    for text, level, bold in content:
        fs = 16 if level == 0 else 13
        if first:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(fs)
            run.font.bold = bold
            run.font.color.rgb = MS_BLUE if bold else DARK
            first = False
        else:
            p = add_paragraph(tf, text, level=level, font_size=fs, bold=bold,
                               color=MS_BLUE if (bold and level == 0) else DARK)

    add_footer_and_slide_number(slide, 7)
    return slide


def build_slide8(prs):
    """Actions & Connectors."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_title(slide, "⚙️ Actions & Connectors")

    tf = add_body_textbox(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.9))

    content = [
        ("📦 BUILT-IN ACTIONS:", 0, True),
        ("Condition (if/else branching), Loop (apply to each, do until)", 1, False),
        ("Compose / Parse JSON, Delay / Wait, Terminate", 1, False),
        ("🤖 AI-POWERED ACTIONS (2025/2026):", 0, True),
        ("Prompt-based AI actions (call GPT-style models)", 1, False),
        ("AI Builder: form processing, entity extraction", 1, False),
        ("Content moderation, summarization", 1, False),
        ("🔌 TOP CONNECTORS:", 0, True),
        ("SharePoint, Outlook, Teams, Dataverse", 1, False),
        ("Azure OpenAI, SQL Server, HTTP/Webhook", 1, False),
        ("Salesforce, ServiceNow, SAP, Power BI", 1, False),
        ("Custom Connectors (any REST/GraphQL API)", 1, False),
    ]

    first = True
    for text, level, bold in content:
        fs = 16 if level == 0 else 13
        if first:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(fs)
            run.font.bold = bold
            run.font.color.rgb = MS_BLUE if bold else DARK
            first = False
        else:
            p = add_paragraph(tf, text, level=level, font_size=fs, bold=bold,
                               color=MS_BLUE if (bold and level == 0) else DARK)

    add_footer_and_slide_number(slide, 8)
    return slide


def build_slide9(prs):
    """How to Create an Agent Flow."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_title(slide, "🛠️ Creating an Agent Flow — 4 Methods")

    tf = add_body_textbox(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.9))

    content = [
        ("METHOD 1: 💬 Natural Language (AI-Assisted)", 0, True),
        ("Describe your automation in plain English", 1, False),
        ("Copilot Studio generates a draft flow automatically", 1, False),
        ("METHOD 2: 🖱️ Visual Designer (Canvas)", 0, True),
        ("Drag & drop triggers, actions, conditions", 1, False),
        ("Connect nodes visually on the canvas", 1, False),
        ("METHOD 3: 📋 Start from a Template", 0, True),
        ("Microsoft provides pre-built flow templates", 1, False),
        ("Customize for your scenario", 1, False),
        ("METHOD 4: 🔄 Convert from Power Automate", 0, True),
        ("Import compatible cloud flows into Copilot Studio", 1, False),
        ("Manage and use them as agent flows", 1, False),
    ]

    first = True
    for text, level, bold in content:
        fs = 16 if level == 0 else 13
        if first:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(fs)
            run.font.bold = bold
            run.font.color.rgb = MS_BLUE if bold else DARK
            first = False
        else:
            p = add_paragraph(tf, text, level=level, font_size=fs, bold=bold,
                               color=MS_BLUE if (bold and level == 0) else DARK)

    add_footer_and_slide_number(slide, 9)
    return slide


def build_slide10(prs):
    """Using Agent Flows as Tools in Agents."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_title(slide, "🤝 Using Agent Flows as Tools in Agents")

    tf = add_body_textbox(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.9))

    content = [
        ('Agent Flows are registered as "TOOLS" the AI agent can call', 0, False),
        ("The agent intelligently selects which tool to use based on user intent", 0, False),
        ("Example Flow:", 0, True),
        ('User: "What\'s the status of ticket #12345?"', 1, False),
        ('Agent detects intent → Calls "GetTicketStatus" tool', 1, False),
        ("Flow connects to ServiceNow API → Returns status", 1, False),
        ('Agent: "Ticket #12345 is Open, assigned to John."', 1, False),
        ("Tool description quality directly impacts AI decision accuracy", 0, False),
        ("Typed inputs (text, number, boolean) and structured outputs", 0, False),
        ("🔑 The agent decides WHEN to call the flow and WHAT parameters to pass — automatically", 0, False),
    ]

    first = True
    for text, level, bold in content:
        fs = 16 if level == 0 else 13
        if first:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(fs)
            run.font.bold = bold
            run.font.color.rgb = DARK
            first = False
        else:
            p = add_paragraph(tf, text, level=level, font_size=fs, bold=bold,
                               color=MS_BLUE if (bold and level == 0) else DARK)

    add_footer_and_slide_number(slide, 10)
    return slide


def build_slide11(prs):
    """Flow Properties & Management."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_title(slide, "📋 Flow Properties & Management")

    tf = add_body_textbox(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.9))

    content = [
        ("Configurable Properties:", 0, True),
        ("Name & Description", 1, False),
        ("Owner (auto-set to creator)", 1, False),
        ("Environment (which Power Platform environment)", 1, False),
        ("Connection References (which accounts/credentials)", 1, False),
        ("Input Parameters (typed: text, number, boolean)", 1, False),
        ("Output Values (returned to the calling agent)", 1, False),
        ("⚠️ Important Limitations:", 0, True),
        ("❌ Cannot be shared or have co-owners in Copilot Studio", 1, False),
        ("❌ Cannot be natively copied between environments", 1, False),
        ("✅ Capabilities:", 0, True),
        ("✅ Built-in test runner available", 1, False),
        ("✅ Run history and logs accessible", 1, False),
        ("✅ Can reference existing Power Platform connections", 1, False),
    ]

    first = True
    for text, level, bold in content:
        fs = 16 if level == 0 else 13
        if first:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(fs)
            run.font.bold = bold
            run.font.color.rgb = MS_BLUE if bold else DARK
            first = False
        else:
            p = add_paragraph(tf, text, level=level, font_size=fs, bold=bold,
                               color=MS_BLUE if (bold and level == 0) else DARK)

    add_footer_and_slide_number(slide, 11)
    return slide


def build_slide12(prs):
    """Autonomous Agents + Agent Flows."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_title(slide, "🚀 Autonomous Agents + Agent Flows")

    add_textbox(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.4),
                '"Always-On AI Automation"',
                font_size=16, bold=True, color=MS_BLUE)

    tf = add_body_textbox(slide, Inches(0.5), Inches(1.65), Inches(12.3), Inches(5.5))

    content = [
        ("Traditional Copilot (Reactive):", 0, True),
        ("User speaks → Agent responds → Done", 1, False),
        ("Autonomous Agent (Proactive):", 0, True),
        ("Event occurs → Agent detects → Decides → Acts", 1, False),
        ("🤖 Autonomous Agent Flow Examples:", 0, True),
        ("Monitor inbox → Triage emails → Route to team", 1, False),
        ("Watch SharePoint → New doc → Summarize with AI", 1, False),
        ("Detect DB anomaly → Alert manager on Teams", 1, False),
        ("Check inventory → Below threshold → Create Purchase Order", 1, False),
        ("🛡️ Responsible AI Guardrails:", 0, True),
        ("Define clear scope & boundaries", 1, False),
        ("Use least-privilege permissions", 1, False),
        ("Keep humans in the loop for critical decisions", 1, False),
        ("Validate all inputs", 1, False),
    ]

    first = True
    for text, level, bold in content:
        fs = 16 if level == 0 else 13
        if first:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(fs)
            run.font.bold = bold
            run.font.color.rgb = MS_BLUE if bold else DARK
            first = False
        else:
            p = add_paragraph(tf, text, level=level, font_size=fs, bold=bold,
                               color=MS_BLUE if (bold and level == 0) else DARK)

    add_footer_and_slide_number(slide, 12)
    return slide


def build_slide13(prs):
    """End-to-End Flow Diagram — IT Help Desk."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_title(slide, "🗺️ End-to-End Agent Flow: IT Help Desk")

    tf = add_body_textbox(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.9))

    content = [
        ('User: "My laptop won\'t connect to Wi-Fi"', 0, False),
        ("Step 1 — Agent Processing:", 0, True),
        ('NLU detects "connectivity issue" intent', 1, False),
        ('Agent decides to call "CreateSupportTicket" tool', 1, False),
        ("Step 2 — Agent Flow Executes:", 0, True),
        ("Extract user info from conversation context", 1, False),
        ("Create ticket in ServiceNow via connector", 1, False),
        ("Send Teams notification to IT channel", 1, False),
        ("Return ticket ID back to agent", 1, False),
        ("Step 3 — Agent Responds:", 0, True),
        ('"Ticket #98765 created! IT will contact you soon."', 1, False),
        ("🔑 The entire loop — from user message to ticket creation — is fully automated", 0, False),
    ]

    first = True
    for text, level, bold in content:
        fs = 16 if level == 0 else 13
        if first:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(fs)
            run.font.bold = bold
            run.font.color.rgb = DARK
            first = False
        else:
            p = add_paragraph(tf, text, level=level, font_size=fs, bold=bold,
                               color=MS_BLUE if (bold and level == 0) else DARK)

    add_footer_and_slide_number(slide, 13)
    return slide


def build_slide14(prs):
    """Key Takeaways."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_title(slide, "💡 Key Takeaways — Day 1")

    tf = add_body_textbox(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.9))

    content = [
        ("1️⃣ Agent Flows are native automations inside Copilot Studio — no separate Power Automate license needed", 0, False),
        ("2️⃣ They follow a Trigger → Action → Connector pattern", 0, False),
        ("3️⃣ Flows can be created via natural language, visual canvas, templates, or import from Power Automate", 0, False),
        ("4️⃣ Flows become TOOLS that AI agents intelligently select and call based on user intent", 0, False),
        ("5️⃣ Autonomous agents use scheduled/event triggers to act proactively without user input", 0, False),
        ("6️⃣ Always apply responsible AI guardrails: scope, least-privilege, human-in-the-loop", 0, False),
        ("📌 Official Docs: learn.microsoft.com/microsoft-copilot-studio/flows-overview", 0, False),
        ("📌 Add Flow to Agent: learn.microsoft.com/microsoft-copilot-studio/flow-agent", 0, False),
        ("📌 Training: learn.microsoft.com/training/modules/use-agent-flows/", 0, False),
    ]

    first = True
    for text, level, bold in content:
        if first:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(16)
            run.font.bold = bold
            run.font.color.rgb = DARK
            first = False
        else:
            add_paragraph(tf, text, level=level, font_size=16, bold=bold)

    add_footer_and_slide_number(slide, 14)
    return slide


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        build_slide1,
        build_slide2,
        build_slide3,
        build_slide4,
        build_slide5,
        build_slide6,
        build_slide7,
        build_slide8,
        build_slide9,
        build_slide10,
        build_slide11,
        build_slide12,
        build_slide13,
        build_slide14,
    ]

    for builder in builders:
        builder(prs)

    output_file = "AI4All_Day1_AgentFlows.pptx"
    prs.save(output_file)
    print(f"✅ Presentation saved: {output_file}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
